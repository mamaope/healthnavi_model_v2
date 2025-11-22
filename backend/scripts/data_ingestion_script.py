import os
import io
import sys
import json
import time
import uuid
import argparse
import logging
import mimetypes
import traceback
import re
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from dotenv import load_dotenv
import boto3
import magic
from tqdm import tqdm
from unstructured.partition.auto import partition
from unstructured.partition.text import partition_text
from pypdf import PdfReader
import openai
from pymilvus import (
    connections,
    FieldSchema,
    CollectionSchema,
    DataType,
    Collection,
    utility,
)
from PIL import Image

# Try to import optional libraries for fallback extraction
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Configure basic logging (will be enhanced by setup_logging() in main())
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[logging.StreamHandler()]
)

# Suppress warnings from PDF parsing libraries
logging.getLogger("fitz").setLevel(logging.ERROR)
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logging.getLogger("pypdf").setLevel(logging.ERROR)
logging.getLogger("unstructured").setLevel(logging.INFO)
logging.getLogger("PIL").setLevel(logging.ERROR)

load_dotenv() 

# Azure OpenAI settings (from .env)
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT") 
AZURE_OPENAI_REGION = os.getenv("AZURE_OPENAI_REGION")
DEPLOYMENT = os.getenv("DEPLOYMENT") 
API_VERSION = os.getenv("API_VERSION", "2024-02-01")

# Milvus/Zilliz (from .env)
MILVUS_URI = os.getenv("MILVUS_URI")
MILVUS_TOKEN = os.getenv("MILVUS_TOKEN")
MILVUS_DB_NAME = os.getenv("MILVUS_DB_NAME", "default")
MILVUS_COLLECTION_NAME = os.getenv("MILVUS_COLLECTION_NAME", "medical_knowledge")

LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO")
EMBEDDING_BATCH_SIZE = int(os.getenv("EMBEDDING_BATCH_SIZE", "64"))

# AWS S3 
AWS_ACCESS_KEY_ID = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_ACCESS_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
AWS_REGION = os.getenv("AWS_REGION", "us-east-1")

# ---- Logging ----
def setup_logging(log_dir: str = "./logs"):
    """Setup logging with both file and console handlers"""
    os.makedirs(log_dir, exist_ok=True)
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    log_file = os.path.join(log_dir, f"ingestion_{timestamp}.log")
    
    file_handler = logging.FileHandler(log_file)
    file_handler.setLevel(getattr(logging, LOG_LEVEL.upper(), logging.INFO))
    file_handler.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s"))
    
    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setLevel(getattr(logging, LOG_LEVEL.upper(), logging.INFO))
    console_handler.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s"))
    
    logger = logging.getLogger("ingestion")
    logger.setLevel(getattr(logging, LOG_LEVEL.upper(), logging.INFO))
    
    # Clear any existing handlers to avoid duplicates
    logger.handlers.clear()
    
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)
    
    # Prevent propagation to root logger to avoid duplicate messages
    logger.propagate = False
    
    return logger, log_file

# Get initial logger (will be replaced by setup_logging in main)
logger = logging.getLogger("ingestion")

# ---- Helpers ----

def init_azure_openai():
    """Configure Azure OpenAI using the new API (>= 1.0.0)"""
    if not AZURE_OPENAI_API_KEY or not AZURE_OPENAI_ENDPOINT or not DEPLOYMENT:
        logger.error("Azure OpenAI config missing in .env. Check AZURE_OPENAI_API_KEY, AZURE_OPENAI_ENDPOINT, DEPLOYMENT.")
        raise RuntimeError("Azure OpenAI configuration missing.")
    logger.info("Azure OpenAI client configured for new API (>= 1.0.0).")

def azure_embed_texts(texts: List[str]) -> List[List[float]]:
    """
    Batch call to Azure OpenAI embeddings using the new API (>= 1.0.0).
    Expects DEPLOYMENT to be set in env.
    Returns list of vectors.
    """
    # guard
    if not texts:
        return []

    from openai import AzureOpenAI
    
    # Create client for each call (or cache it globally if preferred)
    client = AzureOpenAI(
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION,
        azure_endpoint=AZURE_OPENAI_ENDPOINT
    )

    vectors = []
    for i in range(0, len(texts), EMBEDDING_BATCH_SIZE):
        batch = texts[i : i + EMBEDDING_BATCH_SIZE]
        try:
            response = client.embeddings.create(
                input=batch,
                model=DEPLOYMENT 
            )
            # response.data is a list of objects with 'embedding'
            for item in response.data:
                vectors.append(item.embedding)
            time.sleep(0.1)  # small throttle
        except Exception as e:
            logger.exception(f"Azure embedding failed for batch starting at {i}: {e}")
            # fallback: insert zero vectors (or raise) — here we raise to force visibility
            raise
    return vectors

def is_executable_in_path(cmd: str) -> bool:
    """Return True if shell command exists in PATH."""
    from shutil import which
    return which(cmd) is not None

def detect_mime_from_bytes(content: bytes, filename: Optional[str]=None) -> str:
    """
    Return MIME string for given bytes using python-magic fallback to mimetypes.
    """
    try:
        mm = magic.Magic(mime=True)
        mime = mm.from_buffer(content)
        return mime or (mimetypes.guess_type(filename)[0] if filename else "application/octet-stream")
    except Exception:
        return mimetypes.guess_type(filename)[0] if filename else "application/octet-stream"

def read_local_file(path: str) -> bytes:
    with open(path, "rb") as f:
        return f.read()

# ---- Milvus helpers ----

def connect_milvus():
    """
    Connects to Milvus using MILVUS_URI (Zilliz managed). The method depends on how your Zilliz instance authentication is set up.
    We'll assume MILVUS_URI is an https URL to a managed endpoint and MILVUS_TOKEN is an auth token for cloud.
    """
    if not MILVUS_URI:
        raise RuntimeError("MILVUS_URI not set in .env")
    endpoint = MILVUS_URI
    # pymilvus expects host and port; for cloud serverless service, use connections.connect with uri param
    try:
        connections.connect(uri=endpoint, token=MILVUS_TOKEN)
        logger.info(f"Connected to Milvus at {endpoint}")
    except Exception as e:
        logger.exception("Failed to connect to Milvus. Check MILVUS_URI and MILVUS_TOKEN.")
        raise

def ensure_milvus_collection(collection_name: str, dim: int) -> Collection:
    """
    Ensure the Milvus collection exists. If not, create with schema:
    - id (primary key)
    - embedding (FLOAT_VECTOR, dim)
    - payload (JSON string)
    """
    fields = [
        FieldSchema(name="id", dtype=DataType.VARCHAR, max_length=64, is_primary=True, auto_id=False),
        FieldSchema(name="embedding", dtype=DataType.FLOAT_VECTOR, dim=dim),
        FieldSchema(name="payload", dtype=DataType.VARCHAR, max_length=8192),
    ]
    schema = CollectionSchema(fields, description="RAG document chunks")

    # Check if collection exists
    has_col = utility.has_collection(collection_name)
    logger.debug(f"Collection '{collection_name}' exists: {has_col}")
    
    if has_col:
        col = Collection(collection_name)
        col.load()  # Ensure it's loaded
        logger.info(f"Loaded existing collection '{collection_name}' with {col.num_entities} entities.")
        return col
    else:
        logger.info(f"Creating new collection '{collection_name}' with dim={dim}...")
        col = Collection(collection_name, schema=schema)
        logger.info(f"Created collection schema for '{collection_name}'.")
        
        # create an index on embedding vector (HNSW recommended)
        index_params = {
            "index_type": "HNSW",
            "metric_type": "COSINE",
            "params": {"M": 16, "efConstruction": 200}
        }
        col.create_index(field_name="embedding", index_params=index_params)
        logger.info(f"Created HNSW index on collection '{collection_name}'.")
        
        col.load()
        logger.info(f"Loaded new collection '{collection_name}'.")
        
        # Verify collection is actually ready
        if not utility.has_collection(collection_name):
            raise RuntimeError(f"Collection '{collection_name}' was created but cannot be found!")
        
        return col

def get_existing_filenames(collection: Collection) -> set:
    """
    Query Milvus collection to get all existing filenames.
    Returns a set of filenames that are already in the collection.
    Uses pagination to avoid gRPC message size limits.
    """
    try:
        collection.load()
        filenames = set()
        batch_size = 1000 
        offset = 0
        
        logger.info("Retrieving existing filenames from collection (this may take a moment)...")
        
        # gRPC has 4MB message limit - payload field can be large
        # Use very small batches to stay under limit
        batch_size = 100  # Small batch to avoid gRPC message size limit
        last_id = ""
        total_processed = 0
        
        while True:
            # Query records after last_id (lexicographic sorting)
            if last_id:
                expr = f'id > "{last_id}"'
            else:
                expr = "id != ''"
            
            try:
                results = collection.query(
                    expr=expr,
                    output_fields=["id", "payload"],
                    limit=batch_size
                )
            except Exception as e:
                if "larger than max" in str(e) and batch_size > 10:
                    batch_size = max(10, batch_size // 2)
                    logger.debug(f"Reducing batch size to {batch_size} due to gRPC limit")
                    continue
                raise
            
            if not results:
                break
            
            for result in results:
                try:
                    payload = json.loads(result.get("payload", "{}"))
                    filename = payload.get("filename")
                    if filename:
                        filenames.add(filename)
                    # Track last ID for next iteration
                    last_id = result.get("id", last_id)
                except Exception as e:
                    logger.debug(f"Could not parse payload: {e}")
                    continue
            
            total_processed += len(results)
            
            if len(results) < batch_size:
                break
            
            # Log progress every 1000 records
            if total_processed % 1000 == 0:
                logger.info(f"Processed {total_processed} chunks, found {len(filenames)} unique files so far...")
        
        logger.info(f"Found {len(filenames)} unique files in collection (from {total_processed} chunks)")
        return filenames
    except Exception as e:
        logger.warning(f"Could not retrieve existing filenames: {e}. Proceeding without deduplication.")
        return set()

def upsert_to_milvus(collection: Collection, ids: List[str], vectors: List[List[float]], metas: List[Dict[str, Any]]):
    """
    Insert rows into Milvus. metas are JSON-serializable metadata per vector.
    """
    payloads = [json.dumps(m) for m in metas]
    try:
        collection.insert([ids, vectors, payloads])
        collection.flush()  # Ensure data is persisted
        logger.info(f"Inserted {len(ids)} vectors to Milvus collection {collection.name}.")
    except Exception as e:
        logger.exception("Milvus insert failed.")
        raise

# ---- Text extraction & chunking ----

def get_pdf_page_labels_bytes(content: bytes) -> List[str]:
    """Return list of page labels from PDF bytes if pypdf available."""
    try:
        if not PdfReader:
            return []
        reader = PdfReader(io.BytesIO(content))
        n = len(reader.pages)
        # create labels as strings (1-based)
        return [str(i+1) for i in range(n)]
    except Exception:
        return []

def parse_with_unstructured(file_bytes: bytes, filename: str, prefer_hi_res: bool=True) -> Tuple[List[Any], str]:
    """
    Return list of 'elements' and the strategy used.
    Uses unstructured.partition with hi_res then falls back to fast.
    Simpler approach from working script with multiple file type support.
    """
    if partition is None:
        raise RuntimeError("unstructured.partition not available (install 'unstructured' package).")
    file_like = io.BytesIO(file_bytes)
    ext = os.path.splitext(filename)[1].lower()
    is_pdf = ext == ".pdf"
    strategy = "hi_res" if prefer_hi_res else "fast"
    last_exc = None
    
    # Try hi_res first if preferred, otherwise try fast
    for strat in ([strategy] if strategy == "fast" else ["hi_res","fast"]):
        try:
            elems = partition(
                file=file_like,
                strategy=strat,
                languages=["eng"],
                include_page_breaks=True,
                infer_table_structure=True,
                extract_images_in_pdf=is_pdf,
                extract_image_block_types=["Image", "Table", "Figure"]
            )
            logger.info(f"partition(...) succeeded for {filename} using strategy '{strat}'")
            return elems, strat
        except Exception as e:
            last_exc = e
            logger.warning(f"partition strategy '{strat}' failed for {filename}: {e}")
            file_like.seek(0)
            continue
    logger.error(f"All partition strategies failed for {filename}: {last_exc}")
    raise last_exc

def element_text(elem) -> str:
    """Extract text from element - handles both dict and object formats"""
    if isinstance(elem, dict):
        return elem.get("text", "")
    return getattr(elem, "text", "") or ""

def extract_with_pymupdf(file_bytes: bytes) -> List[Dict[str, Any]]:
    """
    Fallback extraction using PyMuPDF for when unstructured fails.
    Returns list of element dicts with text and page metadata.
    """
    if not HAS_PYMUPDF:
        raise RuntimeError("PyMuPDF not available")
    
    elements = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                # Create element-like dict for compatibility
                elements.append({
                    "text": text,
                    "page": page_num + 1,
                    "type": "text"
                })
        doc.close()
        logger.info(f"PyMuPDF extracted text from {len(elements)} pages")
    except Exception as e:
        logger.error(f"PyMuPDF extraction failed: {e}")
        raise
    
    return elements

def extract_with_docx(file_bytes: bytes) -> List[Dict[str, Any]]:
    """
    Fallback extraction using python-docx for DOCX files when unstructured fails.
    Returns list of element dicts with text.
    """
    if not HAS_DOCX:
        raise RuntimeError("python-docx not available")
    
    elements = []
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        full_text_parts = []
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                full_text_parts.append(text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    full_text_parts.append(row_text)
        
        if full_text_parts:
            # Combine into single text and create element
            full_text = "\n\n".join(full_text_parts)
            elements.append({
                "text": full_text,
                "page": None,
                "type": "text"
            })
            logger.info(f"python-docx extracted text from DOCX ({len(full_text)} chars)")
        else:
            logger.warning("python-docx extracted no text from DOCX")
    except Exception as e:
        logger.error(f"python-docx extraction failed: {e}")
        raise
    
    return elements

def extract_with_pptx(file_bytes: bytes) -> List[Dict[str, Any]]:
    """
    Fallback extraction using python-pptx for PPTX files when unstructured fails.
    Returns list of element dicts with text and slide metadata.
    """
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not available")
    
    elements = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
            
            if slide_text_parts:
                slide_text = "\n\n".join(slide_text_parts)
                elements.append({
                    "text": slide_text,
                    "page": slide_num,
                    "type": "text"
                })
        
        if elements:
            logger.info(f"python-pptx extracted text from {len(elements)} slides")
        else:
            logger.warning("python-pptx extracted no text from PPTX")
    except Exception as e:
        logger.error(f"python-pptx extraction failed: {e}")
        raise
    
    return elements

def extract_with_pypdf(file_bytes: bytes) -> List[Dict[str, Any]]:
    """
    Fallback extraction using pypdf for PDF files when unstructured fails.
    Returns list of element dicts with text and page metadata.
    """
    if not PdfReader:
        raise RuntimeError("pypdf not available")
    
    elements = []
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                elements.append({
                    "text": text.strip(),
                    "page": page_num,
                    "type": "text"
                })
        
        if elements:
            logger.info(f"pypdf extracted text from {len(elements)} pages")
        else:
            logger.warning("pypdf extracted no text from PDF")
    except Exception as e:
        logger.error(f"pypdf extraction failed: {e}")
        raise
    
    return elements

def extract_plain_text(file_bytes: bytes, encoding: str = "utf-8") -> List[Dict[str, Any]]:
    """
    Extract text from plain text files.
    Returns list of element dicts with text.
    """
    elements = []
    try:
        # Try UTF-8 first, then fallback to latin-1
        for enc in [encoding, "latin-1", "cp1252", "iso-8859-1"]:
            try:
                text = file_bytes.decode(enc)
                if text.strip():
                    elements.append({
                        "text": text,
                        "page": None,
                        "type": "text"
                    })
                    logger.info(f"Extracted plain text using {enc} encoding")
                    return elements
            except UnicodeDecodeError:
                continue
        
        logger.warning("Could not decode plain text file with any encoding")
    except Exception as e:
        logger.error(f"Plain text extraction failed: {e}")
        raise
    
    return elements

def extract_document_title(file_bytes: bytes, filename: str, elements: List[Any]) -> Optional[str]:
    """
    Extract document title from metadata or content.
    Tries: PDF metadata, first heading element, or first significant text.
    Returns None if no good title found (caller should use filename as fallback).
    """
    ext = os.path.splitext(filename)[1].lower()
    
    # Try PDF metadata title
    if ext == ".pdf":
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            if reader.metadata and reader.metadata.get("/Title"):
                title = reader.metadata.get("/Title").strip()
                if title and 3 < len(title) < 200 and not title.lower().startswith("microsoft"):
                    logger.debug(f"Extracted title from PDF metadata: {title}")
                    return title
        except Exception:
            pass
    
    # Try first heading-like element from unstructured
    if elements:
        for elem in elements[:10]:  # Check first 10 elements only
            text = element_text(elem).strip()
            if not text or len(text) < 3:
                continue
            
            # Check if it's a title/heading element type
            elem_type = type(elem).__name__
            if "Title" in elem_type or "Heading" in elem_type:
                if 3 < len(text) < 200:
                    logger.debug(f"Extracted title from {elem_type}: {text}")
                    return text
            
            # First text that looks like a title (not too long, capitalized)
            if len(text) < 150 and (text.isupper() or text.istitle() or text[0].isupper()):
                # Avoid generic text like "Page 1" or numbers
                if not text.lower().startswith(("page ", "untitled", "document")) and not text.isdigit():
                    logger.debug(f"Extracted title from first text: {text}")
                    return text
    
    # No good title found
    return None

def heading_based_chunking(elements: List[Any], max_chars: int=1000, overlap: int=150) -> List[Dict[str,Any]]:
    """
    Chunk elements using headings as boundaries if possible; otherwise do char-based splits.
    Returns list of dicts: {'text': ..., 'page': ..., 'section': ...}
    Handles both dict format (from PyMuPDF) and object format (from unstructured).
    """
    chunks = []
    buffer = ""
    section = None
    page = None

    def flush_buf():
        nonlocal buffer, section, page
        t = buffer.strip()
        if t:  # Don't filter here, let split_text_to_chunks handle it
            chunks.extend(split_text_to_chunks(t, page=page, section=section, max_chars=max_chars, overlap=overlap))
        buffer = ""

    for elem in elements:
        text = element_text(elem).strip()
        if not text:  # Only skip completely empty text
            continue
        
        # Get page number - handle both dict and object formats
        if isinstance(elem, dict):
            page = elem.get("page", page)
        else:
            # Infer page number metadata if available on element
            page_meta = getattr(elem, "metadata", None)
            try:
                candidate_page = getattr(page_meta, "page_number", None)
                if candidate_page:
                    page = candidate_page
            except Exception:
                pass

        is_heading = False
        if len(text) < 200 and (text.endswith(":") or text.isupper() or text.startswith("Section") or text.startswith("CHAPTER")):
            is_heading = True

        if is_heading:
            # flush current buffer before starting new section
            flush_buf()
            section = text.strip()
            continue

        # append content
        if buffer:
            buffer += "\n\n" + text
        else:
            buffer = text

        # If buffer large, flush
        if len(buffer) > max_chars * 1.5:
            flush_buf()
            # keep section continuous
    # final flush
    flush_buf()
    return chunks

def split_text_to_chunks(text: str, page: Optional[int]=None, section: Optional[str]=None, max_chars: int=800, overlap: int=150) -> List[Dict[str,Any]]:
    """
    Split a long text into overlapping chunks (character-level) with metadata.
    """
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    L = len(text)
    while start < L:
        end = min(L, start + max_chars)
        chunk_text = text[start:end].strip()
        meta = {"page": page, "section": section}
        chunks.append({"text": chunk_text, "meta": meta})
        if end == L:
            break
        start = end - overlap if (end - overlap) > start else end
    return chunks

# ---- High-level processing pipeline ----

class Ingestor:
    def __init__(self, milvus_collection_name: str):
        self.milvus_collection_name = milvus_collection_name
        init_azure_openai()
        connect_milvus()
        # collection will be created at first upsert when we know dim
        self.collection = None

    def process_one(self, file_bytes: bytes, filename: str, source: str) -> List[Dict[str,Any]]:
        """
        Process a single file bytes, return list of chunks with metadata
        """
        mime = detect_mime_from_bytes(file_bytes, filename)
        ext = os.path.splitext(filename)[1].lower()
        logger.info(f"Detected MIME for {filename}: {mime} (extension: {ext})")

        is_pdf = (mime and "pdf" in mime) or ext == ".pdf"
        is_docx = (mime and ("word" in mime or "document" in mime)) or ext in [".docx", ".doc"]
        is_pptx = (mime and ("presentation" in mime or "powerpoint" in mime)) or ext in [".pptx", ".ppt"]
        is_txt = (mime and "text" in mime) or ext == ".txt"

        prefer_hi_res = True
        # If PDF but it's text-based, prefer fast; if likely image/scanned, prefer hi_res
        if is_pdf:
            # quick heuristics to detect scanned pdf (no textual content)
            try:
                if PdfReader:
                    reader = PdfReader(io.BytesIO(file_bytes))
                    if any(page.extract_text() for page in reader.pages if page.extract_text()):
                        # has embedded text -> fast is fine
                        prefer_hi_res = False
            except Exception:
                # if pypdf fails, keep prefer_hi_res True (let unstructured try)
                pass

        elements = None
        extraction_method = None

        # Step 1: Always try unstructured first (best quality)
        try:
            elements, strat = parse_with_unstructured(file_bytes, filename, prefer_hi_res=prefer_hi_res)
            extraction_method = f"unstructured-{strat}"
            logger.info(f"Unstructured extraction succeeded using '{strat}' strategy")
        except Exception as e:
            logger.warning(f"Unstructured partition failed for {filename}: {e}")
            elements = None

        # Step 2: File-type specific fallbacks
        if not elements:
            if is_pdf:
                # PDF fallbacks: PyMuPDF -> pypdf
                if HAS_PYMUPDF:
                    try:
                        logger.info(f"Attempting PyMuPDF extraction for {filename}...")
                        elements = extract_with_pymupdf(file_bytes)
                        if elements:
                            extraction_method = "pymupdf"
                            logger.info(f"PyMuPDF successfully extracted content from {filename}")
                    except Exception as e_pymupdf:
                        logger.warning(f"PyMuPDF extraction failed for {filename}: {e_pymupdf}")
                
                # Try pypdf as last resort for PDFs
                if not elements:
                    try:
                        logger.info(f"Attempting pypdf text extraction for {filename}...")
                        elements = extract_with_pypdf(file_bytes)
                        if elements:
                            # Convert to unstructured-like format
                            full_text = "\n\n".join([elem["text"] for elem in elements])
                            elements = partition_text(text=full_text)
                            extraction_method = "pypdf"
                            logger.info(f"pypdf successfully extracted text from {filename}")
                    except Exception as e_pypdf:
                        logger.warning(f"pypdf extraction failed for {filename}: {e_pypdf}")
            
            elif is_docx:
                # DOCX fallback: python-docx
                if HAS_DOCX:
                    try:
                        logger.info(f"Attempting python-docx extraction for {filename}...")
                        elements = extract_with_docx(file_bytes)
                        if elements:
                            extraction_method = "python-docx"
                            logger.info(f"python-docx successfully extracted content from {filename}")
                    except Exception as e_docx:
                        logger.warning(f"python-docx extraction failed for {filename}: {e_docx}")
            
            elif is_pptx:
                # PPTX fallback: python-pptx
                if HAS_PPTX:
                    try:
                        logger.info(f"Attempting python-pptx extraction for {filename}...")
                        elements = extract_with_pptx(file_bytes)
                        if elements:
                            extraction_method = "python-pptx"
                            logger.info(f"python-pptx successfully extracted content from {filename}")
                    except Exception as e_pptx:
                        logger.warning(f"python-pptx extraction failed for {filename}: {e_pptx}")
            
            elif is_txt:
                # Plain text fallback
                try:
                    logger.info(f"Attempting plain text extraction for {filename}...")
                    elements = extract_plain_text(file_bytes)
                    if elements:
                        extraction_method = "plain-text"
                        logger.info(f"Plain text extraction succeeded for {filename}")
                except Exception as e_txt:
                    logger.warning(f"Plain text extraction failed for {filename}: {e_txt}")

        # If still no elements, raise error
        if not elements:
            error_msg = f"All extraction methods failed for {filename}. File type: {mime}, Extension: {ext}"
            logger.error(error_msg)
            raise RuntimeError(error_msg)

        # Check if elements have actual text content
        total_text_length = sum(len(element_text(elem)) for elem in elements)
        if total_text_length == 0:
            error_msg = f"Extracted 0 characters from {filename} using {extraction_method}"
            logger.error(error_msg)
            raise RuntimeError(error_msg)

        # Extract document title (try to get from doc metadata/content)
        document_title = extract_document_title(file_bytes, filename, elements)
        if not document_title:
            # Fallback to filename without extension
            document_title = os.path.splitext(os.path.basename(filename))[0]
        
        # chunk
        chunks = heading_based_chunking(elements)
        
        # attach top-level metadata
        enriched = []
        for ch in chunks:
            meta = {
                "source": source,
                "filename": os.path.basename(filename),
                "document_title": document_title,
                "ingested_at": datetime.utcnow().isoformat() + "Z",
                "page": ch["meta"].get("page"),
                "section": ch["meta"].get("section"),
                "extraction_method": extraction_method
            }
            text = ch["text"]
            # simple cleaning
            text = " ".join(text.split())
            if not text or len(text) < 30:
                continue
            enriched.append({"text": text, "meta": meta})
        
        if len(enriched) == 0:
            logger.warning(f"File {filename} produced 0 chunks after filtering (had {len(chunks)} raw chunks, {total_text_length} total chars). This may indicate the content is too short or filtered out.")
            # Try to create at least one chunk from all text if filtering removed everything
            all_text = " ".join([element_text(elem) for elem in elements])
            all_text = " ".join(all_text.split())
            if len(all_text) >= 30:
                enriched.append({
                    "text": all_text[:2000],  # Limit to reasonable size
                    "meta": {
                        "source": source,
                        "filename": os.path.basename(filename),
                        "document_title": document_title,
                        "ingested_at": datetime.utcnow().isoformat() + "Z",
                        "page": None,
                        "section": None,
                        "extraction_method": extraction_method
                    }
                })
                logger.info(f"Created single fallback chunk for {filename} ({len(all_text)} chars)")
        
        logger.info(f"From file {filename} produced {len(enriched)} chunks (title: '{document_title}', method: {extraction_method}).")
        return enriched

    def ingest_batch(self, items: List[Dict[str,Any]]):
        """
        items: list of dict with keys: 'text','meta'
        Embeds each text, upserts to Milvus.
        """
        texts = [it["text"] for it in items]
        metas = [it["meta"] for it in items]
        
        # Add the text content to metadata so it's stored in the payload
        enriched_metas = []
        for text, meta in zip(texts, metas):
            enriched_meta = meta.copy()
            enriched_meta["chunk_text"] = text
            # Also add display_page_number field (map from 'page')
            enriched_meta["display_page_number"] = meta.get("page", "?")
            # Add file_path field (map from 'source')
            enriched_meta["file_path"] = meta.get("source", meta.get("filename", "Unknown"))
            enriched_metas.append(enriched_meta)

        vectors = azure_embed_texts(texts)
        dim = len(vectors[0])
        # ensure collection exists (only create once)
        if self.collection is None:
            self.collection = ensure_milvus_collection(self.milvus_collection_name, dim)

        # generate IDs
        ids = [str(uuid.uuid4()) for _ in vectors]
        # upsert
        upsert_to_milvus(self.collection, ids, vectors, enriched_metas)


# ---- File sources: local dir or S3 ----

def list_local_files(input_dir: str) -> List[str]:
    files = []
    for root, _, filenames in os.walk(input_dir):
        for fn in filenames:
            if fn.startswith("."):
                continue
            files.append(os.path.join(root, fn))
    return files

def list_s3_objects(bucket: str, prefix: str="") -> List[Dict[str,str]]:
    s3 = boto3.client("s3",
                      aws_access_key_id=AWS_ACCESS_KEY_ID,
                      aws_secret_access_key=AWS_SECRET_ACCESS_KEY,
                      region_name=AWS_REGION)
    paginator = s3.get_paginator("list_objects_v2")
    out = []
    for page in paginator.paginate(Bucket=bucket, Prefix=prefix):
        for obj in page.get("Contents", []):
            out.append({"Key": obj["Key"], "Size": obj["Size"]})
    return out

def download_s3_object_bytes(bucket: str, key: str) -> bytes:
    s3 = boto3.client("s3",
                      aws_access_key_id=AWS_ACCESS_KEY_ID,
                      aws_secret_access_key=AWS_SECRET_ACCESS_KEY,
                      region_name=AWS_REGION)
    obj = s3.get_object(Bucket=bucket, Key=key)
    return obj["Body"].read()

# ---- CLI / Runner ----

def main():
    parser = argparse.ArgumentParser(description="Robust multi-format document ingestion script")
    parser.add_argument("--refresh", action="store_true", help="Drop existing collection and start fresh")
    parser.add_argument("--input-dir", type=str, help="Local input directory containing docs (optional)")
    parser.add_argument("--s3-bucket", type=str, help="S3 bucket to read docs from (defaults to env AWS_S3_BUCKET)")
    parser.add_argument("--s3-prefix", type=str, default="input/", help="S3 prefix for objects (default: input/)")
    parser.add_argument("--file", type=str, help="Specific file to process (e.g., input/document.pdf)")
    parser.add_argument("--batch-size", type=int, default=32, help="Number of chunks per embedding batch")
    parser.add_argument("--quarantine-dir", type=str, default="./quarantine", help="Where to move failed files (local only)")
    parser.add_argument("--log-dir", type=str, default="./logs", help="Directory for log files")
    parser.add_argument("--dry-run", action="store_true", help="Parse and chunk but do not embed/store")
    parser.add_argument("--continue-on-error", action="store_true", default=True, help="Continue processing even if files fail")
    args = parser.parse_args()

    # Setup logging with file output first
    global logger
    logger, log_file_path = setup_logging(args.log_dir)

    # Default to S3 if no input-dir specified (like old script)
    if not args.input_dir:
        if not args.s3_bucket:
            # Use default from environment
            args.s3_bucket = os.getenv('AWS_S3_BUCKET', 'healthnavi-cdss')
            logger.info(f"No source specified, using S3 bucket from env: {args.s3_bucket}")
    logger.info("="*80)
    logger.info(f"Starting ingestion run at {datetime.utcnow().isoformat()}Z")
    logger.info(f"Log file: {log_file_path}")
    if args.refresh:
        logger.warning("REFRESH MODE: Will drop existing collection and start fresh!")
    logger.info("="*80)

    os.makedirs(args.quarantine_dir, exist_ok=True)

    # Create failed files report
    failed_files_log = os.path.join(args.log_dir, f"failed_files_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.json")
    failed_files = []

    # Handle refresh mode BEFORE creating Ingestor
    if args.refresh:
        # Connect to Milvus first
        from pymilvus import connections, utility
        try:
            # Disconnect any existing connections first to clear cache
            try:
                connections.disconnect("default")
            except:
                pass
            
            connections.connect(uri=MILVUS_URI, token=MILVUS_TOKEN)
            if utility.has_collection(MILVUS_COLLECTION_NAME):
                logger.warning(f"Dropping existing collection '{MILVUS_COLLECTION_NAME}'...")
                utility.drop_collection(MILVUS_COLLECTION_NAME)
                logger.info(f"Dropped collection '{MILVUS_COLLECTION_NAME}'.")
            connections.disconnect("default")
            time.sleep(2)  # Wait for drop to fully complete
        except Exception as e:
            logger.error(f"Error during refresh: {e}")

    ing = Ingestor(MILVUS_COLLECTION_NAME)

    # gather file list
    work_queue = []
    if args.input_dir:
        files = list_local_files(args.input_dir)
        for f in files:
            work_queue.append({"source": "local", "path": f})
    elif args.s3_bucket:
        if args.file:
            # Process specific file only
            work_queue.append({"source": "s3", "path": args.file})
        else:
            # Process all files in bucket/prefix
            objs = list_s3_objects(args.s3_bucket, args.s3_prefix)
            # Filter for supported extensions
            supported_extensions = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.csv', '.html', '.xml', '.json', '.eml', '.msg'}
            for obj in objs:
                ext = os.path.splitext(obj["Key"])[1].lower()
                if ext in supported_extensions:
                    work_queue.append({"source": "s3", "path": obj["Key"]})

    logger.info(f"Found {len(work_queue)} files to process.")
    
    if len(work_queue) == 0:
        logger.warning("No files found to process. Exiting.")
        sys.exit(0)

    # Get existing filenames to skip already processed files
    logger.info("Checking for already processed files in collection...")
    existing_filenames = set()
    try:
        if ing.collection is None:
            # Create collection to check for existing files
            test_dim = 3072  # Azure text-embedding-3-large dimension
            ing.collection = ensure_milvus_collection(MILVUS_COLLECTION_NAME, test_dim)
        existing_filenames = get_existing_filenames(ing.collection)
        logger.info(f"Found {len(existing_filenames)} unique files already in collection")
    except Exception as e:
        logger.warning(f"Could not check for existing files: {e}. Will process all files.")
    
    # Statistics tracking
    stats = {
        "total_files": len(work_queue),
        "successful": 0,
        "failed": 0,
        "skipped": 0,
        "total_chunks": 0,
        "start_time": datetime.utcnow().isoformat()
    }
    
    logger.info(f"Starting to process {len(work_queue)} files...")
    logger.info(f"First file: {work_queue[0]['path']}")
    if len(work_queue) > 1:
        logger.info(f"Last file: {work_queue[-1]['path']}")

    # process one-by-one, but embed in batches
    batch_accum = []
    last_progress_log = time.time()
    last_file_processed = None
    files_processed_set = set()
    
    # Disable tqdm for non-TTY environments (like nohup)
    use_tqdm = sys.stdout.isatty()
    iterator = tqdm(work_queue, desc="Processing files", disable=not use_tqdm) if use_tqdm else work_queue
    
    for idx, entry in enumerate(iterator, 1):
        file_bytes = None
        filename = None
        source = None
        
        # Safety check: ensure we're not stuck in a loop
        current_file = entry.get("path")
        if current_file == last_file_processed:
            logger.error(f"CRITICAL: Attempted to process same file twice in a row: {current_file}")
            logger.error("This indicates a bug in the loop. Skipping to avoid infinite loop.")
            continue
        
        if current_file in files_processed_set:
            logger.warning(f"File already processed: {current_file}. Skipping duplicate.")
            continue
        
        # Log progress every 30 seconds
        current_time = time.time()
        if current_time - last_progress_log >= 30:
            logger.info(f"Progress: {idx}/{stats['total_files']} files | Success: {stats['successful']} | Failed: {stats['failed']} | Chunks: {stats['total_chunks']}")
            logger.info(f"Currently on file: {os.path.basename(current_file)}")
            last_progress_log = current_time
        
        try:
            # Read file
            if entry["source"] == "local":
                path = entry["path"]
                try:
                    file_bytes = read_local_file(path)
                    filename = path
                    source = f"file://{path}"
                except Exception as e:
                    error_msg = f"Failed to read file: {str(e)}"
                    logger.error(f"Cannot read {path}: {error_msg}")
                    failed_files.append({
                        "file": path,
                        "source": "local",
                        "error": error_msg,
                        "timestamp": datetime.utcnow().isoformat()
                    })
                    stats["failed"] += 1
                    if not args.continue_on_error:
                        raise
                    continue
            else:
                key = entry["path"]
                try:
                    file_bytes = download_s3_object_bytes(args.s3_bucket, key)
                    filename = key
                    source = f"s3://{args.s3_bucket}/{key}"
                except Exception as e:
                    error_msg = f"Failed to download from S3: {str(e)}"
                    logger.error(f"Cannot download {key}: {error_msg}")
                    failed_files.append({
                        "file": key,
                        "source": "s3",
                        "bucket": args.s3_bucket,
                        "error": error_msg,
                        "timestamp": datetime.utcnow().isoformat()
                    })
                    stats["failed"] += 1
                    if not args.continue_on_error:
                        raise
                    continue
            
            # Check if file already exists in collection (deduplication)
            base_filename = os.path.basename(filename)
            if base_filename in existing_filenames:
                logger.info(f"[{idx}/{stats['total_files']}] ⏭️  Skipping {base_filename} (already in collection)")
                stats["skipped"] += 1
                continue

            # Process file
            logger.info(f"[{idx}/{stats['total_files']}] Processing {base_filename}...")
            chunks = ing.process_one(file_bytes, filename, source)
            
            # Success - add chunks to batch
            logger.info(f"[{idx}/{stats['total_files']}] ✓ {os.path.basename(filename)}: {len(chunks)} chunks extracted")
            stats["successful"] += 1
            stats["total_chunks"] += len(chunks)
            
            for c in chunks:
                batch_accum.append({"text": c["text"], "meta": c["meta"]})
                if len(batch_accum) >= args.batch_size:
                    if args.dry_run:
                        logger.info(f"[dry-run] would embed {len(batch_accum)} chunks")
                        batch_accum = []
                    else:
                        ing.ingest_batch(batch_accum)
                        batch_accum = []
                        
        except Exception as e:
            logger.exception(f"Failed processing {entry}. Moving to quarantine.")
            stats["failed"] += 1
            failed_files.append({
                "file": entry.get("path", "unknown"),
                "source": entry["source"],
                "error": str(e),
                "timestamp": datetime.utcnow().isoformat()
            })
            # try to save file locally for later analysis
            try:
                qname = os.path.join(args.quarantine_dir, os.path.basename(entry.get("path", "failed")) + "_" + datetime.utcnow().strftime("%Y%m%dT%H%M%SZ"))
                if entry["source"] == "local" and 'path' in entry:
                    import shutil
                    shutil.copy2(entry["path"], qname)
                else:
                    # write the bytes we downloaded (if available)
                    if file_bytes:
                        with open(qname, "wb") as fh:
                            fh.write(file_bytes)
                logger.info(f"Moved failed file to quarantine: {qname}")
            except Exception:
                logger.exception("Quarantine move failed.")
        
        finally:
            # Mark file as processed regardless of success/failure
            if current_file:
                files_processed_set.add(current_file)
                last_file_processed = current_file

    # final partial batch
    if batch_accum:
        if args.dry_run:
            logger.info(f"[dry-run] would embed {len(batch_accum)} chunks (final batch)")
        else:
            ing.ingest_batch(batch_accum)

    # Write failed files report
    if failed_files:
        with open(failed_files_log, "w") as f:
            json.dump(failed_files, f, indent=2)
        logger.warning(f"Failed files report written to: {failed_files_log}")

    # Final statistics
    stats["end_time"] = datetime.utcnow().isoformat()
    logger.info("="*80)
    logger.info("INGESTION SUMMARY")
    logger.info("="*80)
    logger.info(f"Total files:      {stats['total_files']}")
    logger.info(f"Successful:       {stats['successful']}")
    logger.info(f"Skipped:          {stats['skipped']} (already in collection)")
    logger.info(f"Failed:           {stats['failed']}")
    logger.info(f"Total chunks:     {stats['total_chunks']}")
    processed_count = stats['successful'] + stats['skipped']
    logger.info(f"Success rate:     {processed_count/stats['total_files']*100:.1f}%" if stats['total_files'] > 0 else "N/A")
    if failed_files:
        logger.info(f"Failed files log: {failed_files_log}")
    logger.info(f"Log file:         {log_file_path}")
    logger.info("="*80)
    
    if stats['failed'] > 0:
        logger.warning(f"Ingestion completed with {stats['failed']} failures. Check logs for details.")
    else:
        logger.info("Ingestion completed successfully!")
    
    # Exit with appropriate code
    sys.exit(0 if stats['failed'] == 0 else 1)

if __name__ == "__main__":
    main()
